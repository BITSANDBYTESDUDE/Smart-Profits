"""Smart Profits — GSMA MENA Ignite pitch deck (16:9 PPTX)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x11, 0x20)
NAVY2 = RGBColor(0x12, 0x1A, 0x2E)
CYAN = RGBColor(0x4F, 0xD1, 0xC5)
CYAN_D = RGBColor(0x0F, 0x9E, 0x94)
GOLD = RGBColor(0xE8, 0xC5, 0x6B)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
CARD = RGBColor(0x16, 0x21, 0x36)
BLACK = RGBColor(0x0B, 0x11, 0x20)

W = Inches(13.333)
H = Inches(7.5)
ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "screenshots"
LOGO = ROOT.parent.parent / "public" / "brand" / "mark.png"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def set_run(run, size=18, bold=False, color=WHITE, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    fill(sh, color)
    return sh


def add_round(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(sh, color)
    sh.adjustments[0] = 0.08
    return sh


def add_text(slide, l, t, w, h, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)
    return box


def add_paras(slide, l, t, w, h, lines, size=16, color=WHITE, bold_first=False, space=8):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space)
        run = p.add_run()
        run.text = line
        set_run(run, size, bold_first and i == 0, color)
    return box


def footer(slide, page, total=10):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(8), Inches(0.28), "Smart Profits  ·  GSMA MENA Ignite  ·  Theme 4", 11, False, MUTED)
    add_text(slide, Inches(11.4), Inches(7.12), Inches(1.4), Inches(0.28), f"{page:02d} / {total:02d}", 11, False, MUTED, PP_ALIGN.RIGHT)


def chrome(slide, kicker, title):
    add_rect(slide, 0, 0, W, H, NAVY)
    add_rect(slide, 0, 0, W, Inches(0.07), CYAN)
    add_rect(slide, 0, Inches(0.07), W, Inches(0.015), GOLD)
    add_text(slide, Inches(0.55), Inches(0.28), Inches(12), Inches(0.32), kicker.upper(), 12, True, CYAN)
    add_text(slide, Inches(0.55), Inches(0.58), Inches(12.2), Inches(0.7), title, 28, True, WHITE)


def card(slide, l, t, w, h, heading, body, heading_color=GOLD):
    add_round(slide, l, t, w, h, CARD)
    add_text(slide, l + Inches(0.22), t + Inches(0.16), w - Inches(0.4), Inches(0.36), heading, 16, True, heading_color)
    add_text(slide, l + Inches(0.22), t + Inches(0.52), w - Inches(0.4), h - Inches(0.7), body, 14, False, WHITE)


# ---------------------------------------------------------------------------
# 1 Title
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, 0, W, Inches(0.08), CYAN)
add_rect(s, 0, Inches(0.08), W, Inches(0.02), GOLD)
if LOGO.exists():
    s.shapes.add_picture(str(LOGO), Inches(0.55), Inches(1.35), Inches(1.15), Inches(1.15))
add_text(s, Inches(0.55), Inches(2.55), Inches(12), Inches(0.35), "GSMA MENA IGNITE OPEN GATEWAY HACKATHON  ·  PHASE 1", 13, True, CYAN)
add_text(s, Inches(0.55), Inches(2.95), Inches(12.2), Inches(0.9), "Smart Profits", 48, True, WHITE)
add_text(
    s,
    Inches(0.55),
    Inches(3.85),
    Inches(12),
    Inches(0.7),
    "AI-Driven Financial Analytics & Telco-Secured Merchant Platform",
    22,
    False,
    GOLD,
)
add_text(
    s,
    Inches(0.55),
    Inches(4.6),
    Inches(11.5),
    Inches(0.85),
    "The merchant’s CFO in the chat. Bank-grade SIM-swap defence on the login.\nCAMARA APIs via Nokia Network-as-Code. Decisions by an AI Agent — not by SMS.",
    16,
    False,
    MUTED,
)
for i, tag in enumerate(["SIM Swap", "Number Verification", "Location Verification", "AI Agent", "Theme 4"]):
    x = Inches(0.55) + i * Inches(2.15)
    sh = add_round(s, x, Inches(5.65), Inches(2.0), Inches(0.42), RGBColor(0x14, 0x3A, 0x3A))
    add_text(s, x, Inches(5.68), Inches(2.0), Inches(0.38), tag, 12, True, CYAN, PP_ALIGN.CENTER)
add_text(s, Inches(0.55), Inches(6.35), Inches(12), Inches(0.4), "Bitsandbytesdude Software Agency  ·  Gaza / MENA  ·  August 2026", 14, False, MUTED)
footer(s, 1)

# ---------------------------------------------------------------------------
# 2 Problem
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
chrome(s, "1. Problem statement and context", "The P&L lives in Excel. So does the fraud.")
card(
    s,
    Inches(0.5),
    Inches(1.45),
    Inches(4.0),
    Inches(3.55),
    "Financial blindness",
    "Independent MENA merchants run shops from messy Arabic and English Excel, CSV, PDF and photos. Rent, salaries and utilities are missing. Phantom profit looks healthy. Dead stock and wrong prices drain cash. There is no CFO.",
)
card(
    s,
    Inches(4.65),
    Inches(1.45),
    Inches(4.0),
    Inches(3.55),
    "Profit leakage",
    "Weak SKUs, unexamined margins and unplanned OpEx hide the real net profit. A ‘good month’ in the sales sheet can still be a cash-negative month for the store.",
)
card(
    s,
    Inches(8.8),
    Inches(1.45),
    Inches(4.0),
    Inches(3.55),
    "Account takeover",
    "A SIM swap beats SMS OTP. The attacker then opens the advisor, downloads P&L reports, changes prices, or steals the sales files that are the business. A login from another country uploading ‘the monthly file’ can be a hijacked account.",
)
add_round(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.55), CARD)
add_text(s, Inches(0.75), Inches(5.35), Inches(11.8), Inches(1.25),
         "Banks already buy network APIs. The SaaS that holds the merchant’s books usually does not.\nThat is the Theme 4 gap: financial files without telco-grade identity, on the same login that SMS OTP cannot protect.",
         16, False, WHITE)
footer(s, 2)

# ---------------------------------------------------------------------------
# 3 Solution + APIs
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
chrome(s, "2. Solution and Nokia CAMARA APIs", "One platform. Two brains. Three network APIs.")
card(
    s,
    Inches(0.5),
    Inches(1.4),
    Inches(6.05),
    Inches(1.7),
    "Financial AI",
    "Cleans messy Excel/CSV/PDF. Diagnoses store health, profit leaks, inventory and a 30-day plan. Answers questions from the currently open file — in Arabic or English.",
)
card(
    s,
    Inches(6.75),
    Inches(1.4),
    Inches(6.05),
    Inches(1.7),
    "Smart Guard AI Agent",
    "On login, reset, upload, export and price change, the agent calls Nokia Network-as-Code and decides Allow / Step-up / Freeze. SMS is never the source of truth.",
)

headers = ["Nokia / CAMARA API", "When it is called", "What the agent does"]
rows = [
    ("Number Verification", "Login and step-up identity (silent 4G/5G, no SMS)", "Confirm the live network number. Do not trust a spoofed SMS."),
    ("SIM Swap", "Login, password reset, price change, sensitive actions", "Recent swap → freeze the session. SMS OTP is untrusted."),
    ("Location Verification", "Excel/CSV/PDF upload and P&L export", "Must match the usual store / branch area, or the file never lands."),
]
col_w = [Inches(2.7), Inches(5.0), Inches(4.6)]
table = s.shapes.add_table(4, 3, Inches(0.5), Inches(3.3), Inches(12.3), Inches(2.55)).table
for j, h in enumerate(headers):
    table.columns[j].width = col_w[j]
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x0B, 0x11, 0x20)
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            set_run(r, 12, True, CYAN)
for i, row in enumerate(rows, start=1):
    for j, val in enumerate(row):
        cell = table.cell(i, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD if i % 2 else RGBColor(0x14, 0x1C, 0x30)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run(r, 12, j == 0, WHITE if j else CYAN)
add_text(s, Inches(0.55), Inches(6.0), Inches(12.2), Inches(0.85),
         "Testing and prototype validation will utilize Nokia Network-as-Code developer portal simulators for SIM Swap, Location Verification, and Number Verification — then map one-to-one onto production CAMARA endpoints.",
         13, False, MUTED)
footer(s, 3)

# ---------------------------------------------------------------------------
# 4 AI Agent Design
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
chrome(s, "3. AI Agent design", "Smart Guard decides. It does not caption.")
add_text(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.55),
         "The agent is the orchestration brain required by the hackathon: it calls CAMARA tools, combines them with financial context, and executes policy.",
         16, False, MUTED)

steps = [
    ("1", "Sense", "Merchant action: login, reset, file upload, P&L export, price change."),
    ("2", "Query", "Call Number Verification, SIM Swap and Location Verification via Nokia NaC."),
    ("3", "Decide", "Fuse network risk with file sensitivity. Output: Allow · Step-up · Freeze."),
    ("4", "Act", "If freeze: do not store the workbook. Explain in Arabic or English. Wait for network identity."),
]
for i, (n, title, body) in enumerate(steps):
    x = Inches(0.5) + i * Inches(3.2)
    add_round(s, x, Inches(2.05), Inches(3.0), Inches(2.35), CARD)
    add_text(s, x + Inches(0.18), Inches(2.18), Inches(2.6), Inches(0.4), n, 22, True, GOLD)
    add_text(s, x + Inches(0.18), Inches(2.58), Inches(2.6), Inches(0.35), title, 18, True, CYAN)
    add_text(s, x + Inches(0.18), Inches(3.0), Inches(2.6), Inches(1.2), body, 13, False, WHITE)

add_round(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.15), CARD)
add_text(s, Inches(0.75), Inches(4.75), Inches(12), Inches(0.35), "Freeze loop (example)", 16, True, GOLD)
add_text(
    s,
    Inches(0.75),
    Inches(5.15),
    Inches(11.8),
    Inches(1.4),
    "action: upload monthly_sales.xlsx\n→ agent.check(SIM Swap, Number Verification, Location Verification)\n→ SIM Swap = recent  OR  location ≠ store geofence  →  decision = FREEZE\n→ file is not stored  ·  merchant is notified  ·  recover with Number Verification on the trusted line",
    14,
    False,
    WHITE,
)
footer(s, 4)

# ---------------------------------------------------------------------------
# 5 Architecture
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
chrome(s, "4. Technical architecture", "Next.js  +  PostgreSQL  +  Nokia NaC  +  AI Agent")

boxes = [
    (0.5, 1.45, 12.3, 1.15, "Merchant app  ·  Next.js / TypeScript  ·  Arabic RTL + English  ·  Dashboard, advisor, upload, export"),
    (0.5, 2.85, 6.0, 1.35, "Smart Guard AI Agent\nOrchestrates CAMARA tools + financial context\nAllow / Step-up / Freeze"),
    (6.8, 2.85, 6.0, 1.35, "Financial engine\nParse Excel/CSV/PDF · P&L · leaks · Q&A on open file"),
    (0.5, 4.45, 4.0, 1.35, "PostgreSQL\nUsers, files, P&L,\nrisk event log"),
    (4.7, 4.45, 4.0, 1.35, "Nokia Network-as-Code\nSingle developer path to\noperator CAMARA APIs"),
    (8.9, 4.45, 3.9, 1.35, "CAMARA\nSIM Swap · Number Verif.\nLocation Verification"),
]
for l, t, w, h, text in boxes:
    add_round(s, Inches(l), Inches(t), Inches(w), Inches(h), CARD)
    add_text(s, Inches(l + 0.2), Inches(t + 0.18), Inches(w - 0.4), Inches(h - 0.3), text, 14, True, WHITE)

add_text(s, Inches(0.55), Inches(6.0), Inches(12.2), Inches(0.85),
         "The agent is the only path to login step-up, financial file upload and report export. Prototype testing runs on Nokia Network-as-Code developer portal simulators.",
         14, False, MUTED)
footer(s, 5)

# ---------------------------------------------------------------------------
# 6 Business model
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
chrome(s, "5. Business model & monetization", "Try the books. Then pay for the guard.")

plans = [
    ("7-day free trial", "Full product", "Upload, diagnosis, bilingual advisor and Smart Guard demo. No card required to start. Converts into Pro unless cancelled."),
    ("Pro  ·  $49 / month", "The merchant plan", "Ongoing P&L, profit leaks, inventory advice, open-file Q&A, and telco-secured login / upload / export for a single store."),
    ("Business  ·  $99 / month", "Multi-branch", "Several locations, tighter Location Verification geofences, team seats, and priority freeze review."),
]
for i, (title, kicker, body) in enumerate(plans):
    x = Inches(0.5) + i * Inches(4.2)
    add_round(s, x, Inches(1.45), Inches(3.95), Inches(3.35), CARD)
    add_text(s, x + Inches(0.22), Inches(1.62), Inches(3.5), Inches(0.35), kicker.upper(), 11, True, CYAN)
    add_text(s, x + Inches(0.22), Inches(2.0), Inches(3.5), Inches(0.7), title, 20, True, GOLD)
    add_text(s, x + Inches(0.22), Inches(2.75), Inches(3.5), Inches(1.8), body, 14, False, WHITE)

add_round(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.75), CARD)
add_text(s, Inches(0.75), Inches(5.15), Inches(12), Inches(0.35), "Why this monetizes Open Gateway", 16, True, GOLD)
add_text(s, Inches(0.75), Inches(5.55), Inches(11.8), Inches(1.05),
         "SaaS subscription is the merchant revenue. CAMARA calls are the security cost of goods — and a future B2B2X path with operators (usage-based Number Verification / SIM Swap / Location on every sensitive action). Smart Profits is incubated as a product line of Bitsandbytesdude, not a one-hack demo.",
         14, False, WHITE)
footer(s, 6)

# ---------------------------------------------------------------------------
# 7 Demo screenshots
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
chrome(s, "6. Demo screenshots", "A working merchant product — not a slide-only idea")

shot_files = [
    ("02-dashboard.png", "Dashboard — health, P&L, phantom vs real profit"),
    ("03-advisor.png", "Advisor — file-grounded Q&A"),
    ("04-data.png", "Files — messy Excel cleaned automatically"),
    ("05-simulator.png", "Simulator — what-if pricing and forecast"),
]
present = [(SHOTS / name, cap) for name, cap in shot_files if (SHOTS / name).exists()]
if not present:
    present = [(p, p.stem) for p in sorted(SHOTS.glob("*.png"))[:4]]

if present:
    for i, (path, cap) in enumerate(present[:4]):
        col, row = i % 2, i // 2
        x = Inches(0.45) + col * Inches(6.4)
        y = Inches(1.35) + row * Inches(2.7)
        pic_w, pic_h = Inches(6.15), Inches(2.25)
        s.shapes.add_picture(str(path), x, y, pic_w, pic_h)
        add_text(s, x, y + Inches(2.26), Inches(6.15), Inches(0.32), cap, 12, True, CYAN)
else:
    add_text(s, Inches(0.55), Inches(2.2), Inches(12), Inches(1), "Prototype screenshots will be attached from the live Next.js app.", 18, False, MUTED)

footer(s, 7)

# ---------------------------------------------------------------------------
# 8 Demo links + walkthrough
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
chrome(s, "6. Demo walkthrough & links", "What judges will see in three minutes")
walk = [
    ("1  Analyse", "Upload a messy furniture ledger. Ask the advisor for the highest-profit SKU. See phantom vs real net profit."),
    ("2  Attack", "Nokia NaC simulator: SIM Swap = true. Upload and export freeze. The agent explains why SMS OTP is untrusted."),
    ("3  Place", "Location mismatch vs store geofence. The financial file never lands. Number Verification recovers the trusted line."),
]
for i, (title, body) in enumerate(walk):
    y = Inches(1.4) + i * Inches(1.25)
    add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.12), CARD)
    add_text(s, Inches(0.75), y + Inches(0.12), Inches(12), Inches(0.32), title, 18, True, GOLD)
    add_text(s, Inches(0.75), y + Inches(0.48), Inches(11.8), Inches(0.5), body, 15, False, WHITE)

add_text(s, Inches(0.55), Inches(5.3), Inches(12.2), Inches(1.5),
         "Product repository:  https://github.com/BITSANDBYTESDUDE/Smart-Profits\nStudio / company:     https://bitsandbytesdude.vercel.app/\nNetwork layer:          Nokia Network-as-Code developer portal simulators (SIM Swap, Number Verification, Location Verification)",
         15, False, CYAN)
footer(s, 8)

# ---------------------------------------------------------------------------
# 9 Team
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
chrome(s, "7. Team bios and roles", "A Gaza product, inside a shipping software studio")

add_round(s, Inches(0.5), Inches(1.4), Inches(6.05), Inches(3.85), CARD)
add_text(s, Inches(0.75), Inches(1.55), Inches(5.55), Inches(0.3), "FOUNDER  /  PRODUCT LEAD", 12, True, CYAN)
add_text(s, Inches(0.75), Inches(1.9), Inches(5.55), Inches(0.45), "Israa Nael Hamad", 22, True, GOLD)
add_text(s, Inches(0.75), Inches(2.4), Inches(5.55), Inches(0.4), "Software Engineering  ·  final year", 14, True, WHITE)
add_text(
    s,
    Inches(0.75),
    Inches(2.9),
    Inches(5.55),
    Inches(2.1),
    "University of Palestine, Gaza.\nOwns product vision, merchant UX, bilingual (AR/EN) FinTech experience, and the Theme 4 narrative: phantom Excel profit vs SIM-swap takeover of the same files. Smart Profits is her product, brought into the studio as a core line.",
    14,
    False,
    WHITE,
)

add_round(s, Inches(6.75), Inches(1.4), Inches(6.05), Inches(3.85), CARD)
add_text(s, Inches(7.0), Inches(1.55), Inches(5.55), Inches(0.3), "PRINCIPAL  /  TECHNICAL CO-FOUNDER", 12, True, CYAN)
add_text(s, Inches(7.0), Inches(1.9), Inches(5.55), Inches(0.45), "Mir Kashif", 22, True, GOLD)
add_text(s, Inches(7.0), Inches(2.4), Inches(5.55), Inches(0.4), "Bitsandbytesdude Software Agency", 14, True, WHITE)
add_text(
    s,
    Inches(7.0),
    Inches(2.9),
    Inches(5.55),
    Inches(2.1),
    "Leads engineering through the studio: full-stack platforms, AI & agent workflows, and cloud delivery. Smart Profits is being added to the Bitsandbytesdude product portfolio — Next.js architecture, agent orchestration, and Nokia NaC integration under one shipping team.",
    14,
    False,
    WHITE,
)

add_round(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.35), CARD)
add_text(s, Inches(0.75), Inches(5.55), Inches(12), Inches(1.05),
         "Studio: Bitsandbytesdude — high-performance software, AI automation and product systems.\nhttps://bitsandbytesdude.vercel.app/   ·   bitsandbytesdude@gmail.com   ·   Remote-first, MENA-rooted",
         15, False, WHITE)
footer(s, 9)

# ---------------------------------------------------------------------------
# 10 Close
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, 0, W, Inches(0.08), CYAN)
add_rect(s, 0, Inches(0.08), W, Inches(0.02), GOLD)
add_text(s, Inches(0.55), Inches(1.7), Inches(12), Inches(0.35), "THE ASK", 14, True, CYAN)
add_text(s, Inches(0.55), Inches(2.15), Inches(12.2), Inches(1.3), "Advance Smart Profits under Theme 4.", 36, True, WHITE)
add_text(
    s,
    Inches(0.55),
    Inches(3.55),
    Inches(12),
    Inches(1.2),
    "A real merchant SaaS. A real AI Agent. Real CAMARA APIs through Nokia Network-as-Code.\nNext: NaC simulator freeze demo  ·  path to MWC Doha 2026.",
    18,
    False,
    MUTED,
)
for i, tag in enumerate(["CAMARA via Nokia NaC — used", "AI Agent orchestration — used", "MENA inclusion — AR + EN"]):
    x = Inches(0.55) + i * Inches(4.1)
    add_round(s, x, Inches(5.1), Inches(3.9), Inches(0.5), RGBColor(0x14, 0x3A, 0x3A))
    add_text(s, x, Inches(5.15), Inches(3.9), Inches(0.42), tag, 13, True, CYAN, PP_ALIGN.CENTER)
add_text(s, Inches(0.55), Inches(5.9), Inches(12), Inches(0.7),
         "Israa Nael Hamad  ·  Mir Kashif  ·  Bitsandbytesdude  ·  University of Palestine, Gaza",
         14, False, GOLD)
footer(s, 10)

out = ROOT / "Smart-Profits-PITCH-DECK.pptx"
prs.save(out)
print("saved", out)
print("shots used", [p.name for p, _ in present] if present else [])
